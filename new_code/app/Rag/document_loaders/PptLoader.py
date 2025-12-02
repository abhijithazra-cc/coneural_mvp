from langchain_community.document_loaders import PyMuPDFLoader
from Rag.abstractions.Idocloader import Idocloader
import fitz  # PyMuPDF (best)
from io import BytesIO
from langchain_core.documents import Document
from typing import Optional
import docx
from io import BytesIO
import pandas as pd
from pptx import Presentation
from io import BytesIO

def extract_text_from_ppt_bytes(ppt_bytes: bytes):
    prs = Presentation(BytesIO(ppt_bytes))
    full_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)

    return "\n".join(full_text)
class PptLoader(Idocloader):
     def __init__(self):
          self.documents=None

     def load_document(self,file:str | bytes,filename:str):
                
                if type(file)==str:
                    print("file type","str")
                    docs = []
                    prs = Presentation(file)
  

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                            
                               docs.append(Document(page_content=shape.text, metadata={}))
                    self.documents=docs
                    return docs


                else:
                    docs = []
                    prs = Presentation(BytesIO(file))
  

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                            
                               docs.append(Document(page_content=shape.text, metadata={}))
                    self.documents=docs
                    return docs

                 

     def get_document(self):
          return self.documents
     def get_full_content(self):
          res=""
          for item in self.documents:
              res+=" "+item.page_content

          return res
